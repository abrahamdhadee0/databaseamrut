import os
import io
import pandas as pd
import threading
from typing import List, Optional, Dict, Tuple
from fastapi import FastAPI, File, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
from sqlalchemy import create_engine, inspect, MetaData, Table, text

# AI & LangChain
from langchain_groq import ChatGroq
from langchain_experimental.agents.agent_toolkits import create_pandas_dataframe_agent
from langchain_community.agent_toolkits import create_sql_agent
from langchain_community.agent_toolkits import SQLDatabaseToolkit
from langchain_community.utilities import SQLDatabase

# File Reading Libraries
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# Load Environment Variables
load_dotenv()

app = FastAPI()

# Enable CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# GLOBAL STORAGE
stored_files = {
    "dfs": [],
    "df_names": [],
    "document_text": "",
    "db_schema": {},  # Store database schema
    "table_relationships": {}  # Store table relationships
}

# Setup AI Engine
llm = ChatGroq(
    model="llama-3.3-70b-versatile",
    temperature=0,
    groq_api_key=os.getenv("GROQ_API_KEY")
)

# Database Connection
db_url = os.getenv("DATABASE_URL")


# ============================================================
# DATABASE INTELLIGENCE FUNCTIONS
# ============================================================

def analyze_database_schema(engine) -> Dict:
    """Analyze entire database schema and relationships"""
    inspector = inspect(engine)
    metadata = MetaData()
    metadata.reflect(bind=engine)

    schema_info = {
        "tables": {},
        "total_tables": 0,
        "total_columns": 0,
        "relationships": [],
        "common_columns": {}
    }

    table_names = inspector.get_table_names()
    schema_info["total_tables"] = len(table_names)

    print(f"\n🔍 Analyzing database schema ({len(table_names)} tables)...")

    for table_name in table_names:
        columns = inspector.get_columns(table_name)
        pk_constraint = inspector.get_pk_constraint(table_name)
        foreign_keys = inspector.get_foreign_keys(table_name)
        indexes = inspector.get_indexes(table_name)

        # Get sample data to understand data types better
        try:
            sample_query = f"SELECT * FROM {table_name} LIMIT 5"
            sample_df = pd.read_sql(sample_query, engine)
            row_count = pd.read_sql(f"SELECT COUNT(*) as count FROM {table_name}", engine)['count'][0]
        except:
            sample_df = pd.DataFrame()
            row_count = 0

        schema_info["tables"][table_name] = {
            "columns": [
                {
                    "name": col["name"],
                    "type": str(col["type"]),
                    "nullable": col["nullable"],
                    "default": col.get("default")
                }
                for col in columns
            ],
            "primary_key": pk_constraint.get("constrained_columns", []),
            "foreign_keys": [
                {
                    "column": fk["constrained_columns"],
                    "references_table": fk["referred_table"],
                    "references_column": fk["referred_columns"]
                }
                for fk in foreign_keys
            ],
            "indexes": [idx["name"] for idx in indexes],
            "row_count": row_count,
            "sample_data": sample_df.head(3).to_dict() if not sample_df.empty else {}
        }

        schema_info["total_columns"] += len(columns)

        # Track relationships
        for fk in foreign_keys:
            schema_info["relationships"].append({
                "from_table": table_name,
                "from_column": fk["constrained_columns"],
                "to_table": fk["referred_table"],
                "to_column": fk["referred_columns"]
            })

        print(f"   ✅ {table_name}: {len(columns)} columns, {row_count:,} rows")

    # Find common columns across tables (potential join keys)
    all_columns = {}
    for table_name, table_info in schema_info["tables"].items():
        for col in table_info["columns"]:
            col_name = col["name"]
            if col_name not in all_columns:
                all_columns[col_name] = []
            all_columns[col_name].append(table_name)

    # Identify columns that appear in multiple tables
    schema_info["common_columns"] = {
        col_name: tables
        for col_name, tables in all_columns.items()
        if len(tables) > 1
    }

    return schema_info


def identify_status_column(table_info: Dict) -> Optional[str]:
    """Identify which column represents status/active state"""
    status_keywords = ['status', 'active', 'is_active', 'state', 'enabled', 'deleted', 'is_deleted']

    for col in table_info['columns']:
        col_name_lower = col['name'].lower()
        for keyword in status_keywords:
            if keyword in col_name_lower:
                return col['name']

    return None


def identify_name_column(table_info: Dict) -> Optional[str]:
    """Identify which column represents the name"""
    name_keywords = ['name', 'title', 'project_name', 'description']

    for col in table_info['columns']:
        col_name_lower = col['name'].lower()
        for keyword in name_keywords:
            if keyword in col_name_lower:
                return col['name']

    return None


def execute_smart_sql(engine, schema_info: Dict, user_query: str) -> Optional[str]:
    """Execute SQL queries with intelligent column detection and proper formatting"""

    query_lower = user_query.lower()

    # Find the main table being queried
    main_table = None
    for table_name in schema_info['tables'].keys():
        table_simple = table_name.replace("tbl_", "").replace("_", " ").lower()
        if table_simple in query_lower or table_name.lower() in query_lower:
            main_table = table_name
            break

    if not main_table:
        return None

    table_info = schema_info['tables'][main_table]

    # PATTERN 1: Active/Inactive projects with names
    if any(word in query_lower for word in ['active', 'inactive', 'status']):
        status_col = identify_status_column(table_info)
        name_col = identify_name_column(table_info)

        if status_col and name_col:
            try:
                # Get all records with status and name
                sql = f"SELECT {name_col}, {status_col} FROM {main_table}"
                print(f"   🚀 Smart SQL: {sql}")
                df = pd.read_sql(sql, engine)

                if df.empty:
                    return f"No records found in {main_table}."

                # Analyze status column to determine active/inactive
                status_values = df[status_col].unique()
                print(f"   📊 Status values found: {status_values}")

                # Common patterns for active status
                active_indicators = [1, '1', 'active', 'Active', 'ACTIVE', True, 'true', 'True', 'yes', 'Yes']
                inactive_indicators = [0, '0', 'inactive', 'Inactive', 'INACTIVE', False, 'false', 'False', 'no', 'No',
                                       'deleted', 'Deleted']

                # Separate active and inactive
                active_mask = df[status_col].isin(active_indicators)
                active_projects = df[active_mask]
                inactive_projects = df[~active_mask]

                # Format response
                response = f"## 📊 Project Status Summary\n\n"
                response += f"**Total Projects:** {len(df)}\n"
                response += f"**Active Projects:** {len(active_projects)}\n"
                response += f"**Inactive Projects:** {len(inactive_projects)}\n\n"

                if len(active_projects) > 0:
                    response += f"### ✅ Active Projects ({len(active_projects)}):\n"
                    for idx, row in active_projects.iterrows():
                        response += f"{idx + 1}. {row[name_col]}\n"
                    response += "\n"

                if len(inactive_projects) > 0:
                    response += f"### ❌ Inactive Projects ({len(inactive_projects)}):\n"
                    for idx, row in inactive_projects.iterrows():
                        response += f"{idx + 1}. {row[name_col]}\n"

                return response

            except Exception as e:
                print(f"   ❌ Smart SQL failed: {e}")
                return None

    # PATTERN 2: Simple count
    if any(phrase in query_lower for phrase in ["number of", "how many", "count"]):
        try:
            sql = f"SELECT COUNT(*) as count FROM {main_table}"
            print(f"   🚀 Smart SQL: {sql}")
            result = pd.read_sql(sql, engine)
            count = result['count'][0]
            readable_name = main_table.replace("tbl_", "").replace("_", " ").title()
            return f"**Total {readable_name}:** {count:,}"
        except Exception as e:
            print(f"   ❌ Count query failed: {e}")
            return None

    # PATTERN 3: List all with specific columns
    if any(phrase in query_lower for phrase in ["list", "show", "display"]):
        name_col = identify_name_column(table_info)
        if name_col:
            try:
                sql = f"SELECT {name_col} FROM {main_table} LIMIT 20"
                print(f"   🚀 Smart SQL: {sql}")
                df = pd.read_sql(sql, engine)

                response = f"## {main_table.replace('tbl_', '').replace('_', ' ').title()}\n\n"
                for idx, row in df.iterrows():
                    response += f"{idx + 1}. {row[name_col]}\n"

                if len(df) == 20:
                    response += f"\n*(Showing first 20 of {len(df)} records)*"

                return response
            except Exception as e:
                print(f"   ❌ List query failed: {e}")
                return None

    return None


def generate_database_context(schema_info: Dict) -> str:
    """Generate comprehensive database context for the AI"""

    context = f"""
=== DATABASE OVERVIEW ===
Total Tables: {schema_info['total_tables']}
Total Columns: {schema_info['total_columns']}
Total Relationships: {len(schema_info['relationships'])}

=== KEY TABLES (Top 20) ===
"""

    # Show top 20 tables by row count
    sorted_tables = sorted(
        schema_info['tables'].items(),
        key=lambda x: x[1]['row_count'],
        reverse=True
    )[:20]

    for table_name, table_info in sorted_tables:
        context += f"\n{table_name} ({table_info['row_count']:,} rows)\n"
        context += "  Columns: " + ", ".join([col['name'] for col in table_info['columns'][:5]])
        if len(table_info['columns']) > 5:
            context += f" ... ({len(table_info['columns'])} total)"
        context += "\n"

    return context


def create_intelligent_database_agent(engine, schema_info: Dict):
    """Create a database-aware agent with strict accuracy controls"""

    db = SQLDatabase(engine)
    toolkit = SQLDatabaseToolkit(db=db, llm=llm)

    # Generate comprehensive context
    db_context = generate_database_context(schema_info)

    # Get list of main tables
    main_tables = ", ".join(list(schema_info['tables'].keys())[:30])

    prefix = f"""You are a PRECISE database analyst. Your job is to execute EXACT SQL queries and return ONLY TRUE DATA from the database.

{db_context}

CRITICAL RULES - NEVER VIOLATE:
1. **NEVER MAKE UP DATA** - Only return data that actually exists in query results
2. **EXECUTE SQL FIRST** - Always query the database before answering
3. **ONE QUERY RULE** - Execute ONE SQL query, get results, provide Final Answer
4. **NO HALLUCINATION** - If you don't see it in the SQL result, DON'T say it
5. **VERIFY RESULTS** - Double-check that your answer matches the SQL output

AVAILABLE TABLES: {main_tables}

QUERY WORKFLOW:
1. Understand the question
2. Identify the relevant table and columns
3. Write ONE precise SQL query
4. Execute it ONCE
5. Read the ACTUAL results
6. Format the answer using ONLY the data from results
7. Provide "Final Answer: [formatted answer from actual data]"

FORMATTING RULES:
- Use markdown for better readability
- Use bullet points for lists
- Use bold for numbers and key terms
- Organize data clearly

EXAMPLE:
Question: "Show me active and inactive projects with names"
Step 1: Identify table (tbl_project) and columns (project_name, status)
Step 2: SQL Query: SELECT project_name, status FROM tbl_project
Step 3: Results: 
   - Project A, 1
   - Project B, 0
   - Project C, 1
Step 4: Final Answer:
**Active Projects (2):**
- Project A
- Project C

**Inactive Projects (1):**
- Project B

REMEMBER: Use ONLY data from SQL results. NO fabrication allowed!
"""

    agent = create_sql_agent(
        llm=llm,
        toolkit=toolkit,
        verbose=True,
        agent_type="openai-tools",
        max_iterations=6,
        max_execution_time=30,
        early_stopping_method="force",
        handle_parsing_errors=True,
        prefix=prefix,
        agent_executor_kwargs={
            "return_intermediate_steps": False,
            "handle_parsing_errors": True
        }
    )

    return agent


# ============================================================
# FILE READING FUNCTIONS
# ============================================================

def read_pdf(file_bytes):
    """Extracts text from PDF with improved formatting preservation."""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text = ""
        page_count = len(reader.pages)

        for page_num, page in enumerate(reader.pages, 1):
            text += f"\n{'=' * 60}\n"
            text += f"PAGE {page_num} of {page_count}\n"
            text += f"{'=' * 60}\n\n"

            content = page.extract_text()
            if content:
                content = content.replace('\x00', '')
                content = content.replace('\r\n', '\n')
                content = content.replace('\r', '\n')

                lines = [line.strip() for line in content.split('\n')]
                content = '\n'.join(line for line in lines if line)

                text += content + "\n\n"

        if not text.strip() or len(text.strip()) < 50:
            return "[WARNING: This PDF appears to be an image/scan or is empty.]"

        return text
    except Exception as e:
        return f"[Error reading PDF: {e}]"


def read_word(file_bytes):
    """Extracts text from .docx with improved structure preservation"""
    try:
        doc = Document(io.BytesIO(file_bytes))
        text = []

        text.append("=" * 60)
        text.append("WORD DOCUMENT CONTENT")
        text.append("=" * 60 + "\n")

        for para in doc.paragraphs:
            if para.text.strip():
                if para.style.name.startswith('Heading'):
                    text.append(f"\n{'=' * 40}")
                    text.append(f"{para.text.upper()}")
                    text.append(f"{'=' * 40}\n")
                else:
                    text.append(para.text)

        if doc.tables:
            text.append("\n" + "=" * 60)
            text.append("TABLES IN DOCUMENT")
            text.append("=" * 60 + "\n")

            for table_num, table in enumerate(doc.tables, 1):
                text.append(f"\nTable {table_num}:")
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                    if row_text:
                        text.append(row_text)
                text.append("")

        return "\n".join(text)
    except Exception as e:
        return f"[Error reading Word Doc: {e}]"


def read_ppt(file_bytes):
    """Extracts text from .pptx slides with improved structure"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = []

        text.append("=" * 60)
        text.append(f"POWERPOINT PRESENTATION ({len(prs.slides)} slides)")
        text.append("=" * 60 + "\n")

        for slide_num, slide in enumerate(prs.slides, 1):
            text.append(f"\n{'=' * 40}")
            text.append(f"SLIDE {slide_num}")
            text.append(f"{'=' * 40}\n")

            if slide.shapes.title:
                text.append(f"TITLE: {slide.shapes.title.text}")
                text.append("")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape != slide.shapes.title:
                        text.append(shape.text)

            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text.append("\nNOTES:")
                    text.append(notes_text)

            text.append("")

        return "\n".join(text)
    except Exception as e:
        return f"[Error reading PowerPoint: {e}]"


def process_documents_with_llm(document_text: str, user_query: str) -> str:
    """Process documents with OPTIMIZED prompt"""

    query_lower = user_query.lower()
    word_count = len(document_text.split())
    char_count = len(document_text)
    line_count = len([line for line in document_text.split('\n') if line.strip()])
    page_count = document_text.count("PAGE ") if "PAGE " in document_text else "Unknown"

    is_summary = any(
        word in query_lower for word in ["summarize", "summary", "overview", "about", "what is", "explain"])
    is_specific = any(
        word in query_lower for word in ["find", "search", "look for", "extract", "get", "show me", "list"])
    is_analysis = any(word in query_lower for word in ["analyze", "compare", "evaluate", "assess", "review"])

    if is_summary:
        task_instruction = """TASK: Provide a comprehensive summary with key points, data, names, conclusions, and action items."""
    elif is_specific:
        task_instruction = """TASK: Search thoroughly and extract ALL relevant information with exact quotes, numbers, and context."""
    elif is_analysis:
        task_instruction = """TASK: Provide deep analysis with patterns, comparisons, insights, and conclusions."""
    else:
        task_instruction = """TASK: Answer thoroughly with direct answer, evidence, context, and details."""

    prompt = f"""You are AMRUT DATA AI - an advanced document intelligence system.

=== DOCUMENT STATISTICS ===
Pages: {page_count} | Words: {word_count:,} | Lines: {line_count:,}

=== DOCUMENT CONTENT ===
{document_text}

=== USER QUESTION ===
{user_query}

=== YOUR MISSION ===
{task_instruction}

CRITICAL RULES:
1. READ CAREFULLY - Analyze entire document
2. BE COMPREHENSIVE - Extract ALL relevant details
3. BE ACCURATE - Only state what's in document
4. BE SPECIFIC - Include exact numbers, dates, names
5. BE ORGANIZED - Use markdown formatting
6. BE COMPLETE - Cover all sections
7. ADD VALUE - Provide insights

FORMAT YOUR RESPONSE WITH:
- Clear headers (##)
- Bullet points for lists
- Bold for emphasis (**text**)
- Organized sections

NOW PROVIDE YOUR COMPREHENSIVE RESPONSE:"""

    try:
        response = llm.invoke(prompt)
        return response.content if hasattr(response, 'content') else str(response)
    except Exception as e:
        return f"Error: {str(e)}"


# ============================================================
# MAIN DATABASE PROCESSING WITH ACCURACY CONTROLS
# ============================================================

def process_database_query(engine, schema_info: Dict, user_query: str) -> str:
    """Process queries against database with strict accuracy controls"""

    print("\n🧠 Processing database query...")

    # STEP 1: Try smart SQL with intelligent column detection (FASTEST & MOST ACCURATE)
    smart_result = execute_smart_sql(engine, schema_info, user_query)
    if smart_result:
        print("   ✅ Answered with smart SQL (ACCURATE)")
        return smart_result

    # STEP 2: Use agent with strict accuracy controls
    try:
        print("   🤖 Using SQL Agent with accuracy controls...")
        agent = create_intelligent_database_agent(engine, schema_info)

        # Enhanced query with accuracy emphasis
        enhanced_query = f"""{user_query}

CRITICAL: 
- Execute SQL query FIRST
- Use ONLY the actual data from SQL results
- Format the response clearly with markdown
- DO NOT make up any data"""

        # Execute with timeout protection using threading
        result_container = {"output": None, "error": None, "completed": False}

        def run_agent():
            try:
                response = agent.invoke({"input": enhanced_query})
                result_container["output"] = response.get("output", "Could not process query.")
                result_container["completed"] = True
            except Exception as e:
                result_container["error"] = str(e)
                result_container["completed"] = True

        # Start agent in separate thread
        thread = threading.Thread(target=run_agent)
        thread.daemon = True
        thread.start()

        # Wait for completion or timeout
        thread.join(timeout=35)  # 35 seconds timeout

        # Check if thread is still running (timeout occurred)
        if thread.is_alive():
            print("   ⏱️ Agent timeout - using fallback")

            # Fallback: Try to provide basic info
            query_lower = user_query.lower()
            for table_name in schema_info['tables'].keys():
                table_simple = table_name.replace("tbl_", "").replace("_", " ").lower()
                if table_simple in query_lower or table_name.lower() in query_lower:
                    try:
                        count_query = f"SELECT COUNT(*) as count FROM {table_name}"
                        result = pd.read_sql(count_query, engine)
                        return f"⏱️ Query timed out. Found {result['count'][0]:,} records in {table_name}.\n\nPlease try a more specific question."
                    except:
                        pass

            return "⏱️ The query is taking too long. Please try:\n- Being more specific\n- Asking a simpler question\n- Breaking into smaller questions"

        # Check for errors
        if result_container["error"]:
            raise Exception(result_container["error"])

        # Get the output
        output = result_container["output"]

        if not output:
            return "❌ Could not generate a response. Please try rephrasing your question."

        # Clean output
        if "Final Answer:" in output:
            output = output.split("Final Answer:")[-1].strip()

        # Remove any remaining agent artifacts
        output = output.replace("Action:", "").replace("Action Input:", "").strip()

        return output

    except Exception as e:
        print(f"   ❌ Agent error: {e}")

        # Last resort: Provide basic table info
        query_lower = user_query.lower()
        for table_name in schema_info['tables'].keys():
            table_simple = table_name.replace("tbl_", "").replace("_", " ").lower()
            if table_simple in query_lower or table_name.lower() in query_lower:
                try:
                    count_query = f"SELECT COUNT(*) as count FROM {table_name}"
                    result = pd.read_sql(count_query, engine)
                    table_info = schema_info['tables'][table_name]
                    columns = ", ".join([col['name'] for col in table_info['columns'][:8]])

                    return f"""❌ Query execution failed, but here's what I found:

**Table:** {table_name}
**Total Records:** {result['count'][0]:,}
**Key Columns:** {columns}

Please try asking a more specific question about this table."""
                except:
                    pass

        return f"❌ Error: {str(e)}\n\nPlease try:\n- Being more specific about the table or data\n- Using simpler language\n- Checking available tables with 'show me all tables'"


# ============================================================
# MAIN CHAT ENDPOINT
# ============================================================

@app.post("/chat")
async def chat_with_amrut_ai(
        message: str = Form(...),
        files: List[UploadFile] = File([])
):
    global stored_files

    print(f"\n{'=' * 60}")
    print(f"📥 REQUEST: '{message}'")
    print(f"📁 Files uploaded: {len(files)}")
    print(f"{'=' * 60}")

    document_text = stored_files["document_text"]
    db_schema = stored_files.get("db_schema")

    # ---------------------------------------------------------
    # 1. LOAD AND ANALYZE DATABASE
    # ---------------------------------------------------------
    if db_url and not db_schema:
        try:
            print("\n🔗 Connecting to database...")
            engine = create_engine(db_url)

            # Analyze database schema (one-time operation)
            db_schema = analyze_database_schema(engine)
            stored_files["db_schema"] = db_schema

            print(f"\n✅ Database analyzed successfully!")
            print(f"   📊 {db_schema['total_tables']} tables")
            print(f"   📋 {db_schema['total_columns']} columns")
            print(f"   🔗 {len(db_schema['relationships'])} relationships")

        except Exception as e:
            print(f"\n⚠️ Database connection failed: {e}")
            return {"reply": f"Database connection error: {e}"}

    # ---------------------------------------------------------
    # 2. LOAD UPLOADED FILES
    # ---------------------------------------------------------
    if files and len(files) > 0:
        print(f"\n📁 Processing {len(files)} new files...")
        document_text = ""

        for file in files:
            content = await file.read()
            filename = file.filename.lower()
            print(f"\n   📄 {file.filename}")

            if filename.endswith(".pdf"):
                text = read_pdf(content)
                document_text += f"\n=== PDF: {file.filename} ===\n{text}\n"
                print(f"      ✅ Extracted ({len(text)} chars)")

            elif filename.endswith(".docx"):
                text = read_word(content)
                document_text += f"\n=== Word: {file.filename} ===\n{text}\n"
                print(f"      ✅ Extracted ({len(text)} chars)")

            elif filename.endswith(".pptx"):
                text = read_ppt(content)
                document_text += f"\n=== PPT: {file.filename} ===\n{text}\n"
                print(f"      ✅ Extracted ({len(text)} chars)")

        stored_files["document_text"] = document_text

    # ---------------------------------------------------------
    # 3. PROCESS QUERY
    # ---------------------------------------------------------

    has_database = db_schema is not None
    has_documents = len(document_text.strip()) > 0

    print(f"\n🧠 ROUTING:")
    print(f"   Database: {has_database} | Documents: {has_documents}")

    try:
        # CASE 1: Database query
        if has_database and not has_documents:
            print("   → Database Query")
            engine = create_engine(db_url)
            response_text = process_database_query(engine, db_schema, message)
            return {"reply": response_text}

        # CASE 2: Document query
        elif has_documents and not has_database:
            print("   → Document Analysis")
            response_text = process_documents_with_llm(document_text, message)
            return {"reply": response_text}

        # CASE 3: Both database and documents
        elif has_database and has_documents:
            print("   → Hybrid Query")
            engine = create_engine(db_url)

            # Try database first
            db_response = process_database_query(engine, db_schema, message)

            # Then documents
            doc_response = process_documents_with_llm(document_text, message)

            # Combine
            combined_prompt = f"""Combine these answers into one comprehensive response:

Question: {message}

From Database: {db_response}

From Documents: {doc_response}

Provide ONE clear, well-formatted answer using markdown."""

            final_response = llm.invoke(combined_prompt)
            return {"reply": final_response.content if hasattr(final_response, 'content') else str(final_response)}

        # CASE 4: No data
        else:
            print("   → No Data Available")
            if db_url:
                return {
                    "reply": "Database is configured but not loaded yet. Please send a query to initialize the connection."}
            else:
                return {"reply": "Please configure a database or upload files to analyze."}

    except Exception as e:
        print(f"\n🔥 ERROR: {e}")
        import traceback
        traceback.print_exc()
        return {"reply": f"Error: {str(e)}"}


@app.post("/clear")
async def clear_files():
    """Clear all stored data"""
    global stored_files
    stored_files = {
        "dfs": [],
        "df_names": [],
        "document_text": "",
        "db_schema": {},
        "table_relationships": {}
    }
    return {"status": "All data cleared"}


@app.get("/schema")
async def get_database_schema():
    """Get database schema information"""
    db_schema = stored_files.get("db_schema")
    if db_schema:
        return {
            "total_tables": db_schema['total_tables'],
            "total_columns": db_schema['total_columns'],
            "tables": list(db_schema['tables'].keys()),
            "relationships": db_schema['relationships']
        }
    return {"error": "Database not analyzed yet"}


@app.get("/")
async def health_check():
    """Health check endpoint"""
    db_schema = stored_files.get("db_schema")
    return {
        "status": "AMRUT DATA AI - Accuracy Optimized v6.0",
        "database_connected": db_schema is not None,
        "tables_loaded": db_schema['total_tables'] if db_schema else 0,
        "document_loaded": len(stored_files["document_text"]) > 0,
        "features": {
            "smart_sql_with_column_detection": True,
            "strict_accuracy_controls": True,
            "no_hallucination_mode": True,
            "markdown_formatting": True,
            "timeout_protection": True
        }
    }


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
